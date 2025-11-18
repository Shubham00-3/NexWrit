from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from io import BytesIO

class DocGenService:
    def __init__(self):
        pass

    def create_docx(self, project_title: str, sections: list) -> BytesIO:
        """Create a Word document from project data"""
        doc = Document()
        
        # Add title
        title = doc.add_heading(project_title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add sections
        for section in sections:
            # Add section heading
            doc.add_heading(section["title"], 1)
            
            # Add section content
            if section.get("content"):
                doc.add_paragraph(section["content"])
            else:
                doc.add_paragraph("[No content generated yet]")
            
            # Add spacing
            doc.add_paragraph()
        
        # Save to BytesIO
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream

    def create_pptx(self, project_title: str, sections: list) -> BytesIO:
        """Create a PowerPoint presentation from project data"""
        prs = Presentation()
        
        # Set slide width and height (16:9)
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = project_title
        
        # Add content slides
        for section in sections:
            # Use title and content layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = section["title"]
            
            # Set content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            if section.get("content"):
                content = section["content"]
                
                # Parse content into bullet points
                lines = content.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        # Remove existing bullet markers if present
                        if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                            line = line[1:].strip()
                        
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
            else:
                p = text_frame.add_paragraph()
                p.text = "[No content generated yet]"
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream

# Global instance
doc_gen_service = DocGenService()
